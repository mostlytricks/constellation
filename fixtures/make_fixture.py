#!/usr/bin/env python3
"""Generate the synthetic test deck: fixtures/orion-sample.pptx

A deliberately patterned deck so `analyze` has real recurrence to find:
  - 3 content slides share one archetype (title bar + bullets + accent bar)
  - a consistent visual grammar (32pt heads / 18pt body, one accent color)
The .pptx itself is git-ignored; regenerate any time:
    .venv/Scripts/python fixtures/make_fixture.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import sys

# Consoles on Korean Windows default to cp949; printed deck data (em-dashes,
# Hangul, census symbols) must never crash on the console's codepage.
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except (AttributeError, ValueError):
        pass

ACCENT = RGBColor(0x2E, 0x5C, 0xFF)   # the deck's one accent color
BODY = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()  # default template, 4:3

# -- slide 1: title -----------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Project Orion"
s.placeholders[1].text = "A synthetic deck for pattern extraction"

# -- slide 2: agenda ----------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = "Agenda"
tf = s.placeholders[1].text_frame
for item in ("Why we're here", "The three pillars", "Next steps"):
    p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
    p.text = item

# -- slides 3-5: the recurring content archetype ------------------------------
# title bar (32pt bold accent) + body bullets (18pt) + bottom accent bar
CONTENT = [
    ("Pillar One: Extraction", ["Deterministic script", "Facts only, JSON out", "No interpretation in code"]),
    ("Pillar Two: Interpretation", ["Agent reads the dump", "Archetypes cite slide numbers", "OPEN: lines over guesses"]),
    ("Pillar Three: Generalization", ["Recurring archetypes become templates", "Source text stripped", "Design guide extracted"]),
]
for title_text, bullets in CONTENT:
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    title = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.9))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = title_text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(32), True, ACCENT
    run.font.name = "Calibri"

    body = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.5))
    tf = body.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size, run.font.color.rgb = Pt(18), BODY
        run.font.name = "Calibri"

    bar = s.shapes.add_shape(1, Inches(0.5), Inches(6.9), Inches(9.0), Inches(0.15))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

# -- slide 6: closing ---------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "감사합니다"          # Korean on purpose: UTF-8 round-trip check
s.placeholders[1].text = "Questions → the night sky"

out = Path(__file__).parent / "orion-sample.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")

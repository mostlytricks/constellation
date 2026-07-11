#!/usr/bin/env python3
"""Generate the second synthetic test deck: fixtures/meridian-pitch.pptx

A consultant-mood financial pitch ("Project Meridian", all content fictional),
deliberately richer than orion-sample so `analyze` meets the things real decks
are made of:
  - a rewritten THEME (custom scheme colors, Georgia/Arial font scheme) so the
    extractor's `themes` section carries non-default values
  - REAL bullet XML -- buChar squares on content slides, buAutoNum on the
    exec summary, buNone on headings (exercises bullet_census, all kinds)
  - theme-INHERITED placeholders (closing slide) next to explicit styling
  - a TABLE and a CHART (the graphic-frame box kinds deck-spec v0 lacks)
  - recurring archetypes: 3 dividers, 3+1 kicker-headline-content slides,
    2 exhibits -- plus a consistent footer/grid on every content slide
The .pptx itself is git-ignored; regenerate any time:
    .venv/Scripts/python fixtures/make_meridian_fixture.py
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x23, 0x40)    # primary brand
INK = RGBColor(0x3B, 0x46, 0x52)     # body text
GOLD = RGBColor(0xC9, 0xA2, 0x27)    # the one accent
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF4, 0xF7)   # light panel
MUTE = RGBColor(0x8A, 0x94, 0xA0)    # footnotes/footers

FIRM = "Meridian Advisory Partners"
CONFID = f"{FIRM} — Strictly Private & Confidential"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# -- theme rewrite: make the master carry this identity, not Office defaults --
def retheme():
    part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    root = part._element if hasattr(part, "_element") else etree.fromstring(part.blob)
    root.set("name", "Meridian")
    clr = root.find(f"{qn('a:themeElements')}/{qn('a:clrScheme')}")
    clr.set("name", "Meridian")
    for name, val in (("dk2", "0B2340"), ("lt2", "F2F4F7"),
                      ("accent1", "C9A227"), ("accent2", "0B2340")):
        el = clr.find(qn(f"a:{name}"))
        for child in list(el):
            el.remove(child)
        el.append(el.makeelement(qn("a:srgbClr"), {"val": val}))
    fonts = root.find(f"{qn('a:themeElements')}/{qn('a:fontScheme')}")
    fonts.set("name", "Meridian")
    for scheme, face in (("a:majorFont", "Georgia"), ("a:minorFont", "Arial")):
        fonts.find(qn(scheme)).find(qn("a:latin")).set("typeface", face)
    if not hasattr(part, "_element"):
        part._blob = etree.tostring(root, xml_declaration=True,
                                    encoding="UTF-8", standalone=True)


def set_bullet(para, kind, char="▪", num_type="arabicPeriod"):
    """Write real bullet XML (what orion-sample's literal '• ' prefixes lack)."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    if kind == "none":
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    elif kind == "char":
        pPr.set("marL", "228600"), pPr.set("indent", "-228600")
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))
    elif kind == "autonum":
        pPr.set("marL", "285750"), pPr.set("indent", "-285750")
        pPr.append(pPr.makeelement(qn("a:buAutoNum"), {"type": num_type}))


def textbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    return tf


def put(para, text, size, color, font="Arial", bold=False, italic=False):
    run = para.add_run()
    run.text = text
    f = run.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return run


def rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def footer(slide, n):
    tf = textbox(slide, 0.6, 7.1, 6.0, 0.25)
    put(tf.paragraphs[0], CONFID, 8, MUTE)
    tf = textbox(slide, 12.3, 7.1, 0.45, 0.25)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    put(tf.paragraphs[0], str(n), 8, MUTE)


def content_scaffold(slide, kicker, headline, source):
    """The recurring consultant grid: kicker / action headline / gold rule /
    source line. Body zone (y 2.1-6.4") is the caller's."""
    p = textbox(slide, 0.6, 0.45, 12.1, 0.3).paragraphs[0]
    set_bullet(p, "none")
    put(p, kicker.upper(), 11, GOLD, bold=True)
    p = textbox(slide, 0.6, 0.8, 12.1, 0.9).paragraphs[0]
    set_bullet(p, "none")
    put(p, headline, 24, NAVY, font="Georgia", bold=True)
    rect(slide, 0.6, 1.75, 1.2, 0.03, GOLD)
    p = textbox(slide, 0.6, 6.65, 12.1, 0.25).paragraphs[0]
    put(p, source, 9, MUTE, italic=True)


def bullets(slide, items, kind="char", size=14):
    tf = textbox(slide, 0.6, 2.1, 12.1, 4.3)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_bullet(p, kind)
        p.space_after = Pt(10)
        put(p, item, size, INK)


def divider(number, title):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    p = textbox(s, 0.6, 2.4, 3.0, 1.1).paragraphs[0]
    put(p, number, 60, GOLD, font="Georgia", bold=True)
    rect(s, 0.6, 3.55, 1.2, 0.03, GOLD)
    p = textbox(s, 0.6, 3.75, 8.0, 0.8).paragraphs[0]
    put(p, title, 32, PAPER, font="Georgia")
    return s


retheme()

# -- 1 · cover ----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0.6, 3.15, 1.2, 0.03, GOLD)
put(textbox(s, 0.6, 3.3, 12.1, 0.9).paragraphs[0],
    "PROJECT MERIDIAN", 40, PAPER, font="Georgia", bold=True)
put(textbox(s, 0.6, 4.25, 12.1, 0.5).paragraphs[0],
    "Series B Growth Financing — Investor Presentation", 16, PANEL)
put(textbox(s, 0.6, 6.8, 12.1, 0.3).paragraphs[0],
    f"{FIRM.upper()}  ·  STRICTLY PRIVATE & CONFIDENTIAL  ·  JUNE 2026",
    10, GOLD, bold=True)

# -- 2 · executive summary (numbered — buAutoNum) ------------------------------
s = prs.slides.add_slide(BLANK)
content_scaffold(s, "Executive summary",
                 "Meridian is raising $40M to scale a proven unit-economics engine",
                 "Source: Company information. Illustrative synthetic data.")
bullets(s, [
    "Revenue grew from $18M to $41M in two years while gross margin expanded "
    "580bps — growth is compounding, not bought.",
    "Net revenue retention of 128% means the installed base alone funds the "
    "plan; new logos are upside.",
    "EBITDA turns positive in FY27 on current cohort economics — the raise "
    "funds acceleration, not survival.",
    "Proceeds: 55% go-to-market, 30% product & engineering, 15% working capital.",
], kind="autonum", size=15)
footer(s, 2)

# -- 3 · divider --------------------------------------------------------------
divider("01", "Market Context")

# -- 4-5 · the recurring content archetype (buChar bullets) --------------------
s = prs.slides.add_slide(BLANK)
content_scaffold(s, "Market context",
                 "An $18B market growing 12% annually, still under-penetrated",
                 "Source: Meridian analysis of industry reports, 2026. Synthetic.")
bullets(s, [
    "Core segment TAM of $18.4B, expanding at a 12.1% CAGR through 2030.",
    "Mid-market penetration is 14% today versus 61% in the enterprise tier — "
    "the growth is where the incumbents are not.",
    "Regulatory tailwinds: 2025 reporting mandates force replacement of manual "
    "workflows in every target account.",
    "Buying behavior has shifted to platform consolidation — 3.2 point "
    "solutions replaced per deal, up from 1.8 in 2023.",
])
footer(s, 4)

s = prs.slides.add_slide(BLANK)
content_scaffold(s, "Competitive landscape",
                 "Incumbents optimize for scale; Meridian wins on integration",
                 "Source: Meridian analysis; public filings. Synthetic.")
bullets(s, [
    "Two scaled incumbents hold 47% share but carry 15-year-old architectures "
    "— integration projects average 11 months.",
    "Meridian deploys in 6 weeks with a 94% on-time record across 120 "
    "implementations.",
    "Win rate versus the largest incumbent is 58% and rising in competitive "
    "evaluations.",
    "The moat is the data layer: every deployment deepens switching costs.",
])
footer(s, 5)

# -- 6 · divider ---------------------------------------------------------------
divider("02", "Financial Performance")

# -- 7 · exhibit: table --------------------------------------------------------
s = prs.slides.add_slide(BLANK)
content_scaffold(s, "Financial performance",
                 "Revenue triples by FY28 with EBITDA turning positive in FY27",
                 "Source: Company model, June 2026. A = actual, E = estimate. Synthetic.")
ROWS = [
    ("$M", "FY24A", "FY25A", "FY26E", "FY27E", "FY28E"),
    ("Revenue", "18.2", "27.4", "41.0", "61.5", "86.3"),
    ("Gross margin", "58%", "61%", "64%", "66%", "68%"),
    ("EBITDA", "(6.1)", "(4.8)", "(1.9)", "3.2", "9.7"),
]
tbl = s.shapes.add_table(4, 6, Inches(0.6), Inches(2.2),
                         Inches(12.1), Inches(2.4)).table
for r, row in enumerate(ROWS):
    for c, value in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = value
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        font = para.runs[0].font
        font.size, font.name = Pt(12), "Arial"
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = NAVY
            font.color.rgb, font.bold = PAPER, True
        else:
            cell.fill.fore_color.rgb = PAPER if r % 2 else PANEL
            font.color.rgb = INK
            font.bold = c == 0
footer(s, 7)

# -- 8 · exhibit: chart ----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
content_scaffold(s, "Revenue trajectory",
                 "Compounding growth: a 48% CAGR through FY28",
                 "Source: Company model, June 2026. Synthetic.")
chart_data = CategoryChartData()
chart_data.categories = ["FY24A", "FY25A", "FY26E", "FY27E", "FY28E"]
chart_data.add_series("Revenue ($M)", (18.2, 27.4, 41.0, 61.5, 86.3))
chart = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.1),
    Inches(11.5), Inches(4.3), chart_data).chart
chart.has_legend = False   # bars pick up theme accent1 = the gold
footer(s, 8)

# -- 9 · divider -----------------------------------------------------------------
divider("03", "The Ask")

# -- 10 · ask (content archetype + highlight panel) -------------------------------
s = prs.slides.add_slide(BLANK)
content_scaffold(s, "The ask",
                 "$40M Series B to fund three growth levers",
                 "Source: Company information. Synthetic.")
tf = textbox(s, 0.6, 2.1, 12.1, 2.6)
for i, item in enumerate([
    "Go-to-market: double the mid-market sales pods that already return 1.4x "
    "on fully loaded cost in year one.",
    "Product: ship the analytics layer that moves us up-stack into "
    "board-level reporting.",
    "Selective M&A: two tuck-in targets identified in adjacent workflow "
    "segments.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_bullet(p, "char")
    p.space_after = Pt(10)
    put(p, item, 14, INK)
panel = rect(s, 0.6, 5.0, 12.1, 1.3, PANEL)
panel_tf = panel.text_frame
panel_tf.word_wrap = True
p = panel_tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
put(p, "Use of funds — 55% go-to-market  ·  30% product & engineering  ·  "
       "15% working capital", 13, NAVY, bold=True)
footer(s, 10)

# -- 11 · closing (placeholders — styling inherited from the Meridian theme) ------
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Thank you"
s.placeholders[1].text = ("James Park  ·  Managing Director  ·  "
                          "jpark@meridianadvisory.example")
put(textbox(s, 0.6, 6.9, 12.1, 0.4).paragraphs[0],
    "This presentation is illustrative synthetic material generated for "
    "software testing. It describes no real company, transaction, or offer.",
    7.5, MUTE, italic=True)

out = Path(__file__).parent / "meridian-pitch.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")

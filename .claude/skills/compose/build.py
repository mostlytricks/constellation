#!/usr/bin/env python3
"""constellation/compose -- deterministic deck-spec renderer.

Renders a deck-spec JSON to .pptx via python-pptx. No judgment: the
*composing* (idea x template x theme -> spec) is the agent's job, done
before this script runs -- the script only validates and draws. It is
also the seam's enforcement: a spec whose template/theme/box/token
references don't resolve is refused with every error listed.

Shapes are defined in .gravity/deck-spec/SPEC.md (spec_version 1):
text boxes (string / bullet-array fills, real buChar bullets), solid
fills, tables (header + banding tokens), and category charts.

Usage:
    python build.py <deck-spec.json> [-o out.pptx]
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[3]
TEMPLATES = ROOT / "library" / "templates"
THEMES = ROOT / "library" / "themes"
SLIDE_H_EMU = 6858000  # fixed height; width follows the spec's aspect

ALIGNS = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
CHART_TYPES = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED,
               "bar": XL_CHART_TYPE.BAR_CLUSTERED,
               "line": XL_CHART_TYPE.LINE}


def load_json(path, what, errors):
    try:
        return json.loads(Path(path).read_text(encoding="utf-8"))
    except FileNotFoundError:
        errors.append(f"{what} not found: {path}")
    except json.JSONDecodeError as e:
        errors.append(f"{what} is not valid JSON: {path} ({e})")
    return None


def box_kind(box):
    """text | bullets | table | chart | fill-only."""
    kind = box.get("content", "text" if "text_style" in box else "fill-only")
    return kind


def check_fill_value(where, key, box, value, errors):
    kind = box_kind(box)
    if kind == "bullets" and not isinstance(value, list):
        errors.append(f"{where}: box '{key}' takes a bullet array, got "
                      f"{type(value).__name__}")
    elif kind == "text" and not isinstance(value, str):
        errors.append(f"{where}: box '{key}' takes a string, got "
                      f"{type(value).__name__}")
    elif kind == "table":
        rows = value if isinstance(value, list) else []
        if not rows or not all(isinstance(r, list)
                               and all(isinstance(c, str) for c in r)
                               and len(r) == len(rows[0]) for r in rows):
            errors.append(f"{where}: box '{key}' takes an array of equal-length "
                          f"string rows (first row = header)")
    elif kind == "chart":
        cats = value.get("categories") if isinstance(value, dict) else None
        series = value.get("series") if isinstance(value, dict) else None
        if not (isinstance(cats, list) and cats and isinstance(series, list)
                and series
                and all(isinstance(s, dict) and isinstance(s.get("name"), str)
                        and isinstance(s.get("values"), list)
                        and len(s["values"]) == len(cats)
                        and all(isinstance(v, (int, float)) for v in s["values"])
                        for s in series)):
            errors.append(f"{where}: box '{key}' takes "
                          f'{{"categories": [...], "series": [{{"name", "values"}}]}} '
                          f"with values matching categories")


def validate(spec, templates, theme, errors):
    """Every cross-reference in the spec must resolve. Collects, never raises."""
    if spec.get("spec_version") != 1:
        errors.append(f"spec_version must be 1, got {spec.get('spec_version')!r} "
                      f"-- v0 specs need the v1 bump (.gravity/deck-spec/SPEC.md)")
        return
    colors = theme["tokens"]["colors"]
    styles = theme["tokens"]["text_styles"]
    for style_name, style in styles.items():
        if style["color"] not in colors:
            errors.append(f"theme style '{style_name}' names unknown color token "
                          f"'{style['color']}'")

    for tpl in templates.values():
        for box in tpl["boxes"]:
            where = f"template '{tpl['template']}' box '{box['box']}'"
            for key in ("text_style", "label_style"):
                if key in box and box[key] not in styles:
                    errors.append(f"{where}: unknown {key} '{box[key]}'")
            if "fill" in box and box["fill"] not in colors:
                errors.append(f"{where}: unknown color token '{box['fill']}'")
            if "align" in box and box["align"] not in ALIGNS:
                errors.append(f"{where}: align must be one of {sorted(ALIGNS)}")
            if box_kind(box) == "table":
                header = box.get("header", {})
                if header.get("text_style") not in styles:
                    errors.append(f"{where}: table header needs a known text_style")
                if header.get("fill") not in colors:
                    errors.append(f"{where}: table header needs a known fill token")
                for token in box.get("banding", []):
                    if token not in colors:
                        errors.append(f"{where}: unknown banding token '{token}'")
            if box_kind(box) == "chart":
                if box.get("chart", "column") not in CHART_TYPES:
                    errors.append(f"{where}: chart must be one of {sorted(CHART_TYPES)}")
                if "series_fill" in box and box["series_fill"] not in colors:
                    errors.append(f"{where}: unknown series_fill token "
                                  f"'{box['series_fill']}'")

    for i, slide in enumerate(spec["slides"], start=1):
        where = f"slide {slide.get('n', '?')}"
        if slide.get("n") != i:
            errors.append(f"{where}: 'n' must be sequential (expected {i})")
        tpl = templates.get(slide["template"])
        if tpl is None:
            continue  # missing-template error already recorded at load
        if slide["role"] not in tpl["roles"] and not slide.get("stretch"):
            errors.append(f"{where}: role '{slide['role']}' is not in template "
                          f"'{tpl['template']}' roles {tpl['roles']} -- set "
                          f"\"stretch\": true if the stretch is intentional")
        boxes = {b["box"]: b for b in tpl["boxes"]}
        fill = slide.get("fill", {})
        for key, value in fill.items():
            box = boxes.get(key)
            if box is None or box_kind(box) == "fill-only":
                errors.append(f"{where}: fill key '{key}' is not a fillable box "
                              f"of template '{tpl['template']}'")
            else:
                check_fill_value(where, key, box, value, errors)
        for box in tpl["boxes"]:
            if box_kind(box) != "fill-only" and box["box"] not in fill:
                errors.append(f"{where}: box '{box['box']}' has no fill entry "
                              f"-- fill it or write a visible 'OPEN: ...' line")


def emu_box(box, sw, sh):
    return (Emu(int(sw * box["x_pct"] / 100)), Emu(int(sh * box["y_pct"] / 100)),
            Emu(int(sw * box["w_pct"] / 100)), Emu(int(sh * box["h_pct"] / 100)))


def apply_style(run, style, colors):
    run.font.size = Pt(style["size_pt"])
    run.font.bold = style["bold"]
    run.font.italic = style.get("italic", False)
    run.font.color.rgb = RGBColor.from_string(colors[style["color"]])
    if "font" in style:
        run.font.name = style["font"]
    if "font_ea" in style:
        # python-pptx's font.name is latin-only; CJK runs need a:ea or they
        # render on fallback. Appended after font.name so a:latin exists first.
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", style["font_ea"])


def styled_text(value, style):
    return value.upper() if style.get("caps") else value


def set_bullet(para, char):
    """Real buChar with hanging indent -- not a literal prefix."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))


def draw_fill(slide_obj, box, hex_color, sw, sh):
    x, y, w, h = emu_box(box, sw, sh)
    shape = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    shape.line.fill.background()


def draw_text(slide_obj, box, value, style, colors, bullet_char, sw, sh):
    x, y, w, h = emu_box(box, sw, sh)
    tf = slide_obj.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    items = value if isinstance(value, list) else [value]
    for j, item in enumerate(items):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        if box.get("align") in ALIGNS:
            para.alignment = ALIGNS[box["align"]]
        if isinstance(value, list):
            set_bullet(para, bullet_char)
        run = para.add_run()
        run.text = styled_text(item, style)
        apply_style(run, style, colors)


def draw_table(slide_obj, box, rows, styles, colors, sw, sh):
    x, y, w, h = emu_box(box, sw, sh)
    tbl = slide_obj.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    header = box["header"]
    banding = box.get("banding")
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            if r == 0:
                style = styles[header["text_style"]]
            elif c == 0 and "label_style" in box:
                style = styles[box["label_style"]]
            else:
                style = styles[box["text_style"]]
            run = para.add_run()
            run.text = styled_text(value, style)
            apply_style(run, style, colors)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(
                    colors[header["fill"]])
            elif banding:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(
                    colors[banding[(r - 1) % len(banding)]])


def draw_chart(slide_obj, box, data, colors, sw, sh):
    x, y, w, h = emu_box(box, sw, sh)
    chart_data = CategoryChartData()
    chart_data.categories = data["categories"]
    for series in data["series"]:
        chart_data.add_series(series["name"], tuple(series["values"]))
    chart = slide_obj.shapes.add_chart(
        CHART_TYPES[box.get("chart", "column")], x, y, w, h, chart_data).chart
    chart.has_legend = len(data["series"]) > 1
    if "series_fill" in box:
        for plot_series in chart.plots[0].series:
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = RGBColor.from_string(
                colors[box["series_fill"]])


def build(spec, templates, theme, out_path):
    colors = theme["tokens"]["colors"]
    styles = theme["tokens"]["text_styles"]
    bullet_char = theme["tokens"].get("bullet_char", "•")
    prs = Presentation()
    prs.slide_height = Emu(SLIDE_H_EMU)
    prs.slide_width = Emu(int(SLIDE_H_EMU * spec["deck"]["aspect"]))
    sw, sh = prs.slide_width, prs.slide_height
    blank = next((l for l in prs.slide_layouts if l.name == "Blank"),
                 prs.slide_layouts[6])

    for slide in spec["slides"]:
        tpl = templates[slide["template"]]
        slide_obj = prs.slides.add_slide(blank)
        for box in tpl["boxes"]:
            if "fill" in box:
                draw_fill(slide_obj, box, colors[box["fill"]], sw, sh)
            kind = box_kind(box)
            if kind == "fill-only":
                continue
            value = slide["fill"][box["box"]]
            if kind == "table":
                draw_table(slide_obj, box, value, styles, colors, sw, sh)
            elif kind == "chart":
                draw_chart(slide_obj, box, value, colors, sw, sh)
            else:
                draw_text(slide_obj, box, value, styles[box["text_style"]],
                          colors, bullet_char, sw, sh)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def verify(spec, templates, out_path):
    """Round-read the built file: slide count matches, every slide has text,
    every table/chart box came back as a real graphic frame."""
    prs = Presentation(out_path)
    slides = list(prs.slides)
    if len(slides) != len(spec["slides"]):
        sys.exit(f"VERIFY FAILED: built {len(slides)} slides, "
                 f"spec has {len(spec['slides'])}")
    for i, (slide_obj, srow) in enumerate(zip(slides, spec["slides"]), start=1):
        texts = [s.text_frame.text for s in slide_obj.shapes if s.has_text_frame]
        if not any(t.strip() for t in texts):
            sys.exit(f"VERIFY FAILED: slide {i} round-read with no text")
        boxes = templates[srow["template"]]["boxes"]
        want = {"table": sum(1 for b in boxes if box_kind(b) == "table"),
                "chart": sum(1 for b in boxes if box_kind(b) == "chart")}
        got = {"table": sum(1 for s in slide_obj.shapes
                            if getattr(s, "has_table", False)),
               "chart": sum(1 for s in slide_obj.shapes
                            if getattr(s, "has_chart", False))}
        if got != want:
            sys.exit(f"VERIFY FAILED: slide {i} graphic frames {got}, "
                     f"template promises {want}")
    return len(slides)


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("spec", help="path to the deck-spec.json to render")
    ap.add_argument("-o", "--out", help="output .pptx path "
                    "(default: deck.pptx alongside the spec)")
    args = ap.parse_args()

    errors = []
    spec = load_json(args.spec, "deck-spec", errors)
    if errors:
        sys.exit("\n".join(f"ERROR: {e}" for e in errors))

    theme = load_json(THEMES / f"{spec['deck']['theme']}.theme.json", "theme", errors)
    templates = {}
    for name in {s["template"] for s in spec["slides"]}:
        tpl = load_json(TEMPLATES / f"{name}.template.json", f"template '{name}'", errors)
        if tpl:
            templates[name] = tpl
    if not errors:
        validate(spec, templates, theme, errors)
    if errors:
        sys.exit("\n".join(f"ERROR: {e}" for e in sorted(errors)))

    out = Path(args.out) if args.out else Path(args.spec).parent / "deck.pptx"
    build(spec, templates, theme, out)
    n = verify(spec, templates, out)
    stretches = [s["n"] for s in spec["slides"] if s.get("stretch")]
    note = f"  (stretched slides: {stretches})" if stretches else ""
    print(f"wrote {out}  ({n} slides, verified round-read){note}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""constellation/analyze -- deterministic structural extractor.

Facts only: dumps a .pptx's structure to JSON. Pattern *recognition*
(slide archetypes, grid, visual grammar) is the agent's job, done by
reading this dump -- the script never interprets, never guesses.

Usage:
    python extract.py <deck.pptx> [-o structure.json]

Output (UTF-8 JSON):
    deck        -- file name, slide size, aspect ratio
    layouts     -- every layout available per master (the deck's vocabulary)
    themes      -- per master: theme color scheme + font scheme, the master's
                   color map and title/body/other text-style defaults -- what
                   every "(inherit)" below resolves to
    slides[]    -- per slide: layout used + every shape with geometry (% of
                   slide, so decks of different sizes compare), placeholder
                   role, text stats, per-run font/size/color census, per-
                   paragraph bullet census
    aggregates  -- layout usage, shape-type census, global font/size/color/
                   bullet census, shape-count histogram (the raw material
                   for clustering slides into archetypes)
"""
import argparse
import json
import sys
from collections import Counter
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn


def pct(value, total):
    """EMU offset -> % of slide dimension (None-safe: inherited positions)."""
    if value is None or not total:
        return None
    return round(value / total * 100, 1)


def _color_val(elem):
    """First color child of an a: element -> hex string or 'scheme:<name>'."""
    if elem is None:
        return None
    for child in elem:
        tag = etree.QName(child).localname
        if tag == "srgbClr":
            return child.get("val")
        if tag == "sysClr":
            return child.get("lastClr") or child.get("val")
        if tag == "schemeClr":
            return f"scheme:{child.get('val')}"
    return None


def _def_rpr_facts(lvl_el):
    """size/bold/color/font defaults from one a:lvlNpPr, only what's stated."""
    rpr = lvl_el.find(qn("a:defRPr"))
    if rpr is None:
        return {}
    facts = {}
    if rpr.get("sz"):
        facts["size_pt"] = int(rpr.get("sz")) / 100
    if rpr.get("b") == "1":
        facts["bold"] = True
    color = _color_val(rpr.find(qn("a:solidFill")))
    if color:
        facts["color"] = color
    latin = rpr.find(qn("a:latin"))
    if latin is not None and latin.get("typeface"):
        facts["font"] = latin.get("typeface")
    return facts


def theme_facts(master):
    """What '(inherit)' resolves to: the master's theme scheme colors/fonts,
    its color map, and its title/body/other text-style defaults (lvl 1-3).
    Facts only, straight from the OOXML -- absent parts are simply omitted."""
    facts = {}
    try:
        theme_root = etree.fromstring(master.part.part_related_by(RT.THEME).blob)
    except KeyError:
        theme_root = None
    if theme_root is not None:
        facts["theme_name"] = theme_root.get("name")
        clr = theme_root.find(f"{qn('a:themeElements')}/{qn('a:clrScheme')}")
        if clr is not None:
            facts["color_scheme"] = {
                "name": clr.get("name"),
                "colors": {
                    etree.QName(c).localname: _color_val(c) for c in clr
                },
            }
        fonts = theme_root.find(f"{qn('a:themeElements')}/{qn('a:fontScheme')}")
        if fonts is not None:
            def typefaces(tag):
                el = fonts.find(qn(tag))
                if el is None:
                    return {}
                return {
                    script: el.find(qn(f"a:{script}")).get("typeface")
                    for script in ("latin", "ea", "cs")
                    if el.find(qn(f"a:{script}")) is not None
                    and el.find(qn(f"a:{script}")).get("typeface")
                }
            facts["font_scheme"] = {
                "name": fonts.get("name"),
                "major": typefaces("a:majorFont"),
                "minor": typefaces("a:minorFont"),
            }
    clr_map = master.element.find(qn("p:clrMap"))
    if clr_map is not None:
        facts["color_map"] = dict(clr_map.attrib)
    tx_styles = master.element.find(qn("p:txStyles"))
    if tx_styles is not None:
        styles = {}
        for style_tag in ("p:titleStyle", "p:bodyStyle", "p:otherStyle"):
            style_el = tx_styles.find(qn(style_tag))
            if style_el is None:
                continue
            levels = {}
            for n in (1, 2, 3):
                lvl = style_el.find(qn(f"a:lvl{n}pPr"))
                if lvl is not None:
                    lvl_facts = _def_rpr_facts(lvl)
                    if lvl_facts:
                        levels[f"lvl{n}"] = lvl_facts
            if levels:
                styles[style_tag.split(":")[1]] = levels
        if styles:
            facts["master_text_styles"] = styles
    return facts


def bullet_census(text_frame):
    """Per-paragraph bullet facts: explicit char/autonum/none vs '(inherit)'
    (bullet unstated in the paragraph -- resolved by layout/master/defaults)."""
    kinds = Counter()
    for para in text_frame.paragraphs:
        pPr = para._p.pPr
        kind = "(inherit)"
        if pPr is not None:
            for child in pPr:
                tag = etree.QName(child).localname
                if tag == "buNone":
                    kind = "none"
                elif tag == "buChar":
                    kind = f"char:{child.get('char')}"
                elif tag == "buAutoNum":
                    kind = f"autonum:{child.get('type')}"
        kinds[kind] += 1
    return dict(kinds)


def font_census(text_frame):
    """Count fonts/sizes/colors across runs. '(inherit)' = theme-resolved.
    fonts = latin typeface; fonts_ea = East-Asian typeface (a:ea) — CJK runs
    without one render on fallback, so the two censuses diverging is a fact."""
    fonts, fonts_ea, sizes, colors = Counter(), Counter(), Counter(), Counter()
    bold_runs = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            fonts[f.name or "(inherit)"] += 1
            rPr = run._r.rPr
            ea = rPr.find(qn("a:ea")) if rPr is not None else None
            fonts_ea[(ea.get("typeface") if ea is not None else None)
                     or "(inherit)"] += 1
            sizes[f"{f.size.pt:g}" if f.size else "(inherit)"] += 1
            try:
                if f.color and f.color.type is not None:
                    colors[str(f.color.rgb)] += 1
                else:
                    colors["(inherit)"] += 1
            except (AttributeError, TypeError):
                colors["(inherit)"] += 1
            if f.bold:
                bold_runs += 1
    return {
        "fonts": dict(fonts),
        "fonts_ea": dict(fonts_ea),
        "sizes_pt": dict(sizes),
        "colors": dict(colors),
        "bold_runs": bold_runs,
    }


def fill_facts(shape):
    """Solid-fill color if resolvable; anything fancier reported by type only."""
    try:
        fill = shape.fill
        kind = str(fill.type)
        if fill.type is not None and str(fill.type).startswith("SOLID"):
            fc = fill.fore_color
            if fc.type is not None and str(fc.type).startswith("SCHEME"):
                return {"type": kind, "theme_color": str(fc.theme_color)}
            return {"type": kind, "rgb": str(fc.rgb)}
        return {"type": kind}
    except (AttributeError, TypeError, NotImplementedError):
        return None


def shape_info(shape, sw, sh):
    info = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "x_pct": pct(shape.left, sw),
        "y_pct": pct(shape.top, sh),
        "w_pct": pct(shape.width, sw),
        "h_pct": pct(shape.height, sh),
    }
    if shape.is_placeholder:
        phf = shape.placeholder_format
        info["placeholder"] = {"idx": phf.idx, "type": str(phf.type)}
    if shape.has_text_frame:
        text = shape.text_frame.text
        info["text_len"] = len(text)
        info["text_preview"] = text[:80]
        info["para_count"] = len(shape.text_frame.paragraphs)
        info["font_census"] = font_census(shape.text_frame)
        info["bullet_census"] = bullet_census(shape.text_frame)
    fill = fill_facts(shape)
    if fill:
        info["fill"] = fill
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = [shape_info(s, sw, sh) for s in shape.shapes]
    if getattr(shape, "has_table", False):
        tbl = shape.table
        info["table"] = {"rows": len(tbl.rows), "cols": len(tbl.columns)}
    if getattr(shape, "has_chart", False):
        info["chart"] = {"type": str(shape.chart.chart_type)}
    return info


def extract(pptx_path):
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height

    slides = []
    layout_usage = Counter()
    type_census = Counter()
    all_fonts, all_fonts_ea = Counter(), Counter()
    all_sizes, all_colors = Counter(), Counter()
    all_bullets = Counter()
    shape_counts = Counter()

    for i, slide in enumerate(prs.slides, start=1):
        layout_usage[slide.slide_layout.name] += 1
        shapes = [shape_info(s, sw, sh) for s in slide.shapes]
        shape_counts[len(shapes)] += 1
        for s in shapes:
            type_census[s["type"]] += 1
            fc = s.get("font_census")
            if fc:
                all_fonts.update(fc["fonts"])
                all_fonts_ea.update(fc["fonts_ea"])
                all_sizes.update(fc["sizes_pt"])
                all_colors.update(fc["colors"])
            bc = s.get("bullet_census")
            if bc:
                all_bullets.update(bc)
        slides.append({
            "n": i,
            "layout": slide.slide_layout.name,
            "shape_count": len(shapes),
            "shapes": shapes,
        })

    return {
        "deck": {
            "file": Path(pptx_path).name,
            "slide_w_emu": sw,
            "slide_h_emu": sh,
            "aspect": round(sw / sh, 3) if sw and sh else None,
            "slide_count": len(slides),
        },
        "layouts": {
            f"master_{m}": [layout.name for layout in master.slide_layouts]
            for m, master in enumerate(prs.slide_masters, start=1)
        },
        "themes": {
            f"master_{m}": theme_facts(master)
            for m, master in enumerate(prs.slide_masters, start=1)
        },
        "slides": slides,
        "aggregates": {
            "layout_usage": dict(layout_usage),
            "shape_type_census": dict(type_census),
            "font_census": dict(all_fonts),
            "font_ea_census": dict(all_fonts_ea),
            "size_census_pt": dict(all_sizes),
            "color_census": dict(all_colors),
            "bullet_census": dict(all_bullets),
            "shape_count_histogram": {str(k): v for k, v in sorted(shape_counts.items())},
        },
    }


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("pptx", help="path to the .pptx to extract")
    ap.add_argument("-o", "--out", help="output JSON path (default: stdout)")
    args = ap.parse_args()

    if not Path(args.pptx).is_file():
        sys.exit(f"ERROR: not a file: {args.pptx}")

    data = extract(args.pptx)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    if args.out:
        out = Path(args.out)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(text, encoding="utf-8")
        print(f"wrote {out}  ({data['deck']['slide_count']} slides)")
    else:
        print(text)


if __name__ == "__main__":
    main()
